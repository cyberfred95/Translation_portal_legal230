"""
Classes de base pour les tests d'extraction de texte.
"""

from io import BytesIO
from django.test import TestCase
from django.core.files.uploadedfile import InMemoryUploadedFile


class FileProcessingTestCase(TestCase):
    """Classe de base pour les tests d'extraction de texte."""
    
    def setUp(self):
        """Préparation des tests."""
        self.api_key = "test_api_key"
    
    def create_test_file(self, content: bytes, filename: str) -> InMemoryUploadedFile:
        """
        Crée un fichier en mémoire pour les tests.
        
        Args:
            content: Contenu binaire du fichier
            filename: Nom du fichier avec extension
            
        Returns:
            InMemoryUploadedFile: Fichier en mémoire pour les tests
        """
        file_obj = BytesIO(content)
        return InMemoryUploadedFile(
            file_obj,
            None,
            filename,
            'application/octet-stream',
            len(content),
            None
        )
    
    def create_word_file(self, paragraphs: list) -> InMemoryUploadedFile:
        """
        Crée un fichier Word de test.
        
        Args:
            paragraphs: Liste de textes de paragraphes
            
        Returns:
            InMemoryUploadedFile: Fichier Word en mémoire
        """
        try:
            from docx import Document
            
            doc = Document()
            for paragraph_text in paragraphs:
                doc.add_paragraph(paragraph_text)
            
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            return self.create_test_file(buffer.read(), "test.docx")
        except ImportError:
            self.skipTest("python-docx n'est pas installé")
    
    def create_excel_file(self, data: dict) -> InMemoryUploadedFile:
        """
        Crée un fichier Excel de test.
        
        Args:
            data: Dictionnaire {cell: value} (ex: {'A1': 'Cell 1'})
            
        Returns:
            InMemoryUploadedFile: Fichier Excel en mémoire
        """
        try:
            from openpyxl import Workbook
            
            wb = Workbook()
            ws = wb.active
            for cell, value in data.items():
                ws[cell] = value
            
            buffer = BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            
            return self.create_test_file(buffer.read(), "test.xlsx")
        except ImportError:
            self.skipTest("openpyxl n'est pas installé")
    
    def create_powerpoint_file(self, title: str = "Test Title") -> InMemoryUploadedFile:
        """
        Crée un fichier PowerPoint de test.
        
        Args:
            title: Titre de la première diapositive
            
        Returns:
            InMemoryUploadedFile: Fichier PowerPoint en mémoire
        """
        try:
            from pptx import Presentation
            
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            return self.create_test_file(buffer.read(), "test.pptx")
        except ImportError:
            self.skipTest("python-pptx n'est pas installé")
    
    def assert_result_format(self, result_dict: dict):
        """
        Vérifie que le format de résultat est compatible avec Custom_mt.
        
        Args:
            result_dict: Dictionnaire de résultat à vérifier
        """
        self.assertIn("texts", result_dict)
        self.assertIsInstance(result_dict["texts"], list)
        if result_dict["texts"]:
            self.assertIn("text", result_dict["texts"][0])
            self.assertIsInstance(result_dict["texts"][0]["text"], str)
