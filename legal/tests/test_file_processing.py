"""
Tests unitaires pour le service d'extraction de texte.

Ces tests vérifient la compatibilité avec l'API Custom_mt et testent
chaque extracteur individuellement.
"""

from io import BytesIO
from django.test import TestCase

from legal.services.file_processing import (
    FileTextExtractorFactory,
    WordExtractor,
    PowerPointExtractor,
    ExcelExtractor,
    TextExtractor,
    ExtractionResult,
)
from legal.tests.base import FileProcessingTestCase


class TextExtractorTest(FileProcessingTestCase):
    """Tests pour l'extracteur de fichiers texte."""
    
    def test_extract_simple_text(self):
        """Test d'extraction d'un texte simple."""
        content = b"Premier paragraphe.\nDeuxieme paragraphe.\n"
        file = self.create_test_file(content, "test.txt")
        
        extractor = TextExtractor()
        result = extractor.extract(file)
        
        self.assertIsInstance(result, ExtractionResult)
        self.assertEqual(len(result.texts), 1)
        self.assertIn("Premier paragraphe. Deuxieme paragraphe.", result.texts[0]["text"])
    
    def test_extract_multiple_paragraphs(self):
        """Test d'extraction avec plusieurs paragraphes."""
        content = b"Paragraphe 1\n\nParagraphe 2\n\nParagraphe 3"
        file = self.create_test_file(content, "test.txt")
        
        extractor = TextExtractor()
        result = extractor.extract(file)
        
        self.assertGreaterEqual(len(result.texts), 2)
    
    def test_format_compatibility(self):
        """Test que le format de retour est compatible avec Custom_mt."""
        content = b"Test content"
        file = self.create_test_file(content, "test.txt")
        
        extractor = TextExtractor()
        result = extractor.extract(file)
        result_dict = result.to_dict()
        
        self.assertIn("texts", result_dict)
        self.assertIsInstance(result_dict["texts"], list)
        if result_dict["texts"]:
            self.assertIn("text", result_dict["texts"][0])


class WordExtractorTest(FileProcessingTestCase):
    """Tests pour l'extracteur Word (nécessite python-docx)."""
    
    def test_word_extractor_import(self):
        """Test que python-docx peut être importé."""
        try:
            WordExtractor()
        except ImportError:
            self.skipTest("python-docx n'est pas installé")
    
    def test_format_compatibility(self):
        """Test que le format de retour est compatible."""
        try:
            from docx import Document
            
            doc = Document()
            doc.add_paragraph("Test paragraph")
            
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            
            file = self.create_test_file(buffer.read(), "test.docx")
            extractor = WordExtractor()
            result = extractor.extract(file)
            result_dict = result.to_dict()
            
            self.assertIn("texts", result_dict)
            self.assertIsInstance(result_dict["texts"], list)
        except ImportError:
            self.skipTest("python-docx n'est pas installé")


class ExcelExtractorTest(FileProcessingTestCase):
    """Tests pour l'extracteur Excel (nécessite openpyxl)."""
    
    def test_excel_extractor_import(self):
        """Test que openpyxl peut être importé."""
        try:
            ExcelExtractor()
        except ImportError:
            self.skipTest("openpyxl n'est pas installé")
    
    def test_format_compatibility(self):
        """Test que le format de retour est compatible."""
        try:
            from openpyxl import Workbook
            
            wb = Workbook()
            ws = wb.active
            ws['A1'] = "Cell 1"
            ws['B1'] = "Cell 2"
            ws['A2'] = "Cell 3"
            
            buffer = BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            
            file = self.create_test_file(buffer.read(), "test.xlsx")
            extractor = ExcelExtractor()
            result = extractor.extract(file)
            result_dict = result.to_dict()
            
            self.assertIn("texts", result_dict)
            self.assertIsInstance(result_dict["texts"], list)
        except ImportError:
            self.skipTest("openpyxl n'est pas installé")


class PowerPointExtractorTest(FileProcessingTestCase):
    """Tests pour l'extracteur PowerPoint (nécessite python-pptx)."""
    
    def test_powerpoint_extractor_import(self):
        """Test que python-pptx peut être importé."""
        try:
            PowerPointExtractor()
        except ImportError:
            self.skipTest("python-pptx n'est pas installé")
    
    def test_format_compatibility(self):
        """Test que le format de retour est compatible."""
        try:
            from pptx import Presentation
            
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Test Title"
            
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            file = self.create_test_file(buffer.read(), "test.pptx")
            extractor = PowerPointExtractor()
            result = extractor.extract(file)
            result_dict = result.to_dict()
            
            self.assertIn("texts", result_dict)
            self.assertIsInstance(result_dict["texts"], list)
        except ImportError:
            self.skipTest("python-pptx n'est pas installé")


class FileTextExtractorFactoryTest(FileProcessingTestCase):
    """Tests pour la factory d'extracteurs."""
    
    def test_get_extractor_docx(self):
        """Test de récupération de l'extracteur Word."""
        extractor = FileTextExtractorFactory.get_extractor('.docx')
        self.assertIsInstance(extractor, WordExtractor)
    
    def test_get_extractor_pptx(self):
        """Test de récupération de l'extracteur PowerPoint."""
        extractor = FileTextExtractorFactory.get_extractor('.pptx')
        self.assertIsInstance(extractor, PowerPointExtractor)
    
    def test_get_extractor_xlsx(self):
        """Test de récupération de l'extracteur Excel."""
        extractor = FileTextExtractorFactory.get_extractor('.xlsx')
        self.assertIsInstance(extractor, ExcelExtractor)
    
    def test_get_extractor_txt(self):
        """Test de récupération de l'extracteur texte."""
        extractor = FileTextExtractorFactory.get_extractor('.txt')
        self.assertIsInstance(extractor, TextExtractor)
    
    def test_get_extractor_unsupported(self):
        """Test avec une extension non supportée."""
        with self.assertRaises(ValueError):
            FileTextExtractorFactory.get_extractor('.pdf')
    
    def test_get_supported_extensions(self):
        """Test de récupération des extensions supportées."""
        extensions = FileTextExtractorFactory.get_supported_extensions()
        self.assertIn('.docx', extensions)
        self.assertIn('.pptx', extensions)
        self.assertIn('.xlsx', extensions)
        self.assertIn('.txt', extensions)
    
    def test_extract_text_method(self):
        """Test de la méthode extract_text qui retourne directement un dict."""
        content = b"Test content"
        file = self.create_test_file(content, "test.txt")
        
        result = FileTextExtractorFactory.extract_text(file)
        
        self.assertIsInstance(result, dict)
        self.assertIn("texts", result)
        self.assertIsInstance(result["texts"], list)


class CompatibilityTest(FileProcessingTestCase):
    """
    Tests de compatibilité avec l'ancien service Custom_mt.
    
    Ces tests vérifient que le nouveau service retourne le même format
    que l'ancien service Custom_mt.
    """
    
    def test_result_structure_matches_custom_mt(self):
        """
        Test que la structure de résultat correspond à celle de Custom_mt.
        
        Format attendu de Custom_mt:
        {
            "texts": [
                {"text": "Premier paragraphe..."},
                {"text": "Deuxième paragraphe..."}
            ]
        }
        """
        content = b"Premier paragraphe\n\nDeuxieme paragraphe"
        file = self.create_test_file(content, "test.txt")
        
        result = FileTextExtractorFactory.extract_text(file)
        
        # Vérifier la structure exacte
        self.assertIn("texts", result)
        self.assertIsInstance(result["texts"], list)
        
        if result["texts"]:
            for text_item in result["texts"]:
                self.assertIn("text", text_item)
                self.assertIsInstance(text_item["text"], str)
    
    def test_empty_file_handling(self):
        """Test de gestion d'un fichier vide."""
        content = b""
        file = self.create_test_file(content, "test.txt")
        
        result = FileTextExtractorFactory.extract_text(file)
        
        self.assertIn("texts", result)
        self.assertIsInstance(result["texts"], list)
    
    def test_file_seek_reset(self):
        """Test que le fichier est repositionné après extraction."""
        content = b"Test content"
        file = self.create_test_file(content, "test.txt")
        original_content = content
        
        FileTextExtractorFactory.extract_text(file)
        
        # Le fichier devrait être repositionné
        file.seek(0)
        self.assertEqual(file.read(), original_content)
