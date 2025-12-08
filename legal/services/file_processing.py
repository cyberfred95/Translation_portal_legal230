"""
Service d'extraction de texte depuis les fichiers Office.

Ce module remplace l'appel au service Custom_mt pour l'extraction de texte
en utilisant des bibliothèques Python locales.
"""

import os
from abc import ABC, abstractmethod
from dataclasses import dataclass
from typing import List, Dict
from io import BytesIO
from contextlib import contextmanager

from django.core.files.uploadedfile import InMemoryUploadedFile


@dataclass(frozen=True)
class ExtractionResult:
    """Résultat de l'extraction de texte au format compatible avec Custom_mt."""
    texts: List[Dict[str, str]]
    
    def to_dict(self) -> dict:
        """Convertit le résultat en dictionnaire JSON compatible."""
        return {"texts": self.texts}


class FileTextExtractor(ABC):
    """Interface abstraite pour les extracteurs de texte."""
    
    @abstractmethod
    def extract(self, file: InMemoryUploadedFile) -> ExtractionResult:
        """
        Extrait le texte d'un fichier.
        
        Args:
            file: Fichier en mémoire à traiter
            
        Returns:
            ExtractionResult: Résultat au format standardisé
        """
        pass
    
    @contextmanager
    def _reset_file_position(self, file: InMemoryUploadedFile):
        """Context manager pour s'assurer que le fichier est repositionné."""
        try:
            file.seek(0)
            yield
        finally:
            file.seek(0)
    
    def _read_file_content(self, file: InMemoryUploadedFile) -> bytes:
        """Lit le contenu du fichier et le repositionne."""
        file.seek(0)
        content = file.read()
        file.seek(0)
        return content
    
    @staticmethod
    def _create_text_item(text: str) -> Dict[str, str]:
        """Crée un élément de texte au format standardisé."""
        return {"text": text}
    
    @staticmethod
    def _check_dependency(package_name: str):
        """
        Vérifie qu'une dépendance est disponible.
        
        Args:
            package_name: Nom du package à vérifier (ex: 'python-docx')
            
        Raises:
            ImportError: Si le package n'est pas installé
        """
        # Convertir le nom du package pour l'import (python-docx -> docx)
        import_name = package_name.replace('python-', '').replace('-', '_')
        try:
            __import__(import_name)
        except ImportError:
            raise ImportError(
                f"{package_name} n'est pas installé. "
                f"Installez-le avec: pip install {package_name}"
            )


class WordExtractor(FileTextExtractor):
    """Extracteur de texte pour les fichiers Word (.docx)."""
    
    def __init__(self):
        self._check_dependency('python-docx')
        from docx import Document
        self.Document = Document
    
    def extract(self, file: InMemoryUploadedFile) -> ExtractionResult:
        """Extrait le texte d'un document Word."""
        with self._reset_file_position(file):
            content = self._read_file_content(file)
            doc = self.Document(BytesIO(content))
            texts = self._extract_paragraphs(doc) + self._extract_tables(doc)
        
        return ExtractionResult(texts=texts)
    
    def _extract_paragraphs(self, doc) -> List[Dict[str, str]]:
        """Extrait le texte des paragraphes."""
        return [
            self._create_text_item(paragraph.text.strip())
            for paragraph in doc.paragraphs
            if paragraph.text.strip()
        ]
    
    def _extract_tables(self, doc) -> List[Dict[str, str]]:
        """Extrait le texte des tableaux."""
        texts = []
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        texts.append(self._create_text_item(cell_text))
        return texts


class PowerPointExtractor(FileTextExtractor):
    """Extracteur de texte pour les fichiers PowerPoint (.pptx)."""
    
    def __init__(self):
        self._check_dependency('python-pptx')
        from pptx import Presentation
        self.Presentation = Presentation
    
    def extract(self, file: InMemoryUploadedFile) -> ExtractionResult:
        """Extrait le texte d'une présentation PowerPoint."""
        with self._reset_file_position(file):
            content = self._read_file_content(file)
            prs = self.Presentation(BytesIO(content))
            texts = self._extract_slides(prs)
        
        return ExtractionResult(texts=texts)
    
    def _extract_slides(self, prs) -> List[Dict[str, str]]:
        """Extrait le texte de toutes les diapositives."""
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(self._create_text_item(text))
        return texts


class ExcelExtractor(FileTextExtractor):
    """Extracteur de texte pour les fichiers Excel (.xlsx)."""
    
    def __init__(self):
        self._check_dependency('openpyxl')
        from openpyxl import load_workbook
        self.load_workbook = load_workbook
    
    def extract(self, file: InMemoryUploadedFile) -> ExtractionResult:
        """Extrait le texte d'un classeur Excel."""
        with self._reset_file_position(file):
            content = self._read_file_content(file)
            workbook = self.load_workbook(BytesIO(content), data_only=True)
            texts = self._extract_sheets(workbook)
        
        return ExtractionResult(texts=texts)
    
    def _extract_sheets(self, workbook) -> List[Dict[str, str]]:
        """Extrait le texte de toutes les feuilles."""
        texts = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_texts = self._extract_row_values(row)
                if row_texts:
                    texts.append(self._create_text_item(" ".join(row_texts)))
        return texts
    
    def _extract_row_values(self, row) -> List[str]:
        """Extrait les valeurs non vides d'une ligne."""
        return [
            str(cell_value).strip()
            for cell_value in row
            if cell_value is not None and str(cell_value).strip()
        ]


class TextExtractor(FileTextExtractor):
    """Extracteur de texte pour les fichiers texte (.txt)."""
    
    # Encodages à essayer dans l'ordre de préférence
    _ENCODINGS = ('utf-8', 'latin-1')
    
    def extract(self, file: InMemoryUploadedFile) -> ExtractionResult:
        """Extrait le texte d'un fichier texte."""
        with self._reset_file_position(file):
            content = self._read_file_content(file)
            text = self._decode_content(content)
            texts = self._parse_paragraphs(text)
        
        return ExtractionResult(texts=texts)
    
    def _decode_content(self, content: bytes) -> str:
        """Décode le contenu avec plusieurs encodages en fallback."""
        for encoding in self._ENCODINGS:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        # Dernier recours : ignorer les erreurs
        return content.decode('utf-8', errors='ignore')
    
    def _parse_paragraphs(self, text: str) -> List[Dict[str, str]]:
        """Parse le texte en paragraphes séparés par des lignes vides."""
        texts = []
        current_paragraph = []
        
        for line in text.split('\n'):
            stripped_line = line.strip()
            if stripped_line:
                current_paragraph.append(stripped_line)
            elif current_paragraph:
                # Ligne vide = fin de paragraphe
                texts.append(self._create_text_item(" ".join(current_paragraph)))
                current_paragraph = []
        
        # Ajouter le dernier paragraphe s'il existe
        if current_paragraph:
            texts.append(self._create_text_item(" ".join(current_paragraph)))
        
        return texts


class FileTextExtractorFactory:
    """Factory pour créer le bon extracteur selon l'extension du fichier."""
    
    # Mapping des extensions vers les classes d'extracteurs
    _EXTRACTOR_MAP = {
        '.docx': WordExtractor,
        '.pptx': PowerPointExtractor,
        '.xlsx': ExcelExtractor,
        '.txt': TextExtractor,
    }
    
    @classmethod
    def get_supported_extensions(cls) -> List[str]:
        """Retourne la liste des extensions supportées."""
        return list(cls._EXTRACTOR_MAP.keys())
    
    @classmethod
    def get_extractor(cls, file_extension: str) -> FileTextExtractor:
        """
        Retourne l'extracteur approprié pour l'extension donnée.
        
        Args:
            file_extension: Extension du fichier (ex: '.docx')
            
        Returns:
            FileTextExtractor: Instance de l'extracteur approprié
            
        Raises:
            ValueError: Si l'extension n'est pas supportée
        """
        file_extension = file_extension.lower()
        extractor_class = cls._EXTRACTOR_MAP.get(file_extension)
        
        if not extractor_class:
            supported = ', '.join(cls.get_supported_extensions())
            raise ValueError(
                f"Extension '{file_extension}' non supportée. "
                f"Extensions supportées: {supported}"
            )
        
        return extractor_class()
    
    @classmethod
    def extract_text(cls, file: InMemoryUploadedFile) -> dict:
        """
        Extrait le texte d'un fichier et retourne le format compatible Custom_mt.
        
        Args:
            file: Fichier en mémoire à traiter
            
        Returns:
            dict: Résultat au format {"texts": [{"text": "..."}]}
            
        Raises:
            ValueError: Si l'extension n'est pas supportée
        """
        file_extension = os.path.splitext(file.name)[1].lower()
        extractor = cls.get_extractor(file_extension)
        result = extractor.extract(file)
        return result.to_dict()
